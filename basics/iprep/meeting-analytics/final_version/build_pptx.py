"""Build meeting_analytics.pptx — 16 slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from PIL import Image
import io

def clean_img(path):
    """Strip metadata and convert to clean RGB PNG in memory."""
    img = Image.open(path).convert('RGB')
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf

# ── Paths ──────────────────────────────────────────────────────────────────
BASE        = Path(__file__).parent
CHARTS      = BASE / 'outputs_notebook'
SCHEMA_IMG  = Path.home() / 'Documents' / 'meeting_analytics.png'
OUT         = BASE / 'meeting_analytics.pptx'

# ── Colours ────────────────────────────────────────────────────────────────
DARK        = RGBColor(0x22, 0x22, 0x22)   # near-black for titles
MID         = RGBColor(0x44, 0x44, 0x44)   # body text
ACCENT      = RGBColor(0xD6, 0x27, 0x28)   # red accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG    = RGBColor(0xFA, 0xFA, 0xFA)
PILL_BG     = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────
def bg(slide, color=SLIDE_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def title_box(slide, title, subtitle=None):
    txbox(slide, title, 0.5, 0.35, 12.3, 1.0, size=36, bold=True, color=DARK)
    if subtitle:
        txbox(slide, subtitle, 0.5, 1.4, 12.3, 0.6, size=18, color=MID)

def section_slide(slide, title, audience, bullets):
    """Standard layout: title top-left, audience pill, bullet points below."""
    bg(slide)
    # Title
    txbox(slide, title, 0.5, 0.3, 11.5, 0.8, size=26, bold=True, color=DARK)
    # Audience pill
    pill = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(1.05), Inches(2.4), Inches(0.32)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = PILL_BG
    pill.line.fill.background()
    pt = pill.text_frame.paragraphs[0]
    pt.alignment = PP_ALIGN.CENTER
    r = pt.add_run()
    r.text = f'For: {audience}'
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE
    # Bullets
    for i, b in enumerate(bullets):
        txbox(slide, f'• {b}', 0.5, 1.55 + i * 0.6, 12.3, 0.55, size=16, color=MID)

def chart_slide(slide, title, audience, bullets, img_path):
    """Left: title + audience + bullets. Right: chart image."""
    bg(slide)
    # Title
    txbox(slide, title, 0.4, 0.25, 5.6, 0.7, size=20, bold=True, color=DARK)
    # Audience pill
    pill = slide.shapes.add_shape(
        1, Inches(0.4), Inches(0.92), Inches(2.2), Inches(0.28)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = PILL_BG
    pill.line.fill.background()
    pt = pill.text_frame.paragraphs[0]
    pt.alignment = PP_ALIGN.CENTER
    r = pt.add_run()
    r.text = f'For: {audience}'
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = WHITE
    # Bullets
    for i, b in enumerate(bullets):
        txbox(slide, f'• {b}', 0.4, 1.3 + i * 0.72, 5.6, 0.65, size=13, color=MID)
    # Chart image
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(
            clean_img(img_path), Inches(6.2), Inches(0.2), Inches(6.9), Inches(7.0)
        )

def bullets_only_slide(slide, title, audience, bullets):
    """Slide with no chart image — bullets centred."""
    bg(slide)
    txbox(slide, title, 0.5, 0.3, 12.3, 0.8, size=26, bold=True, color=DARK)
    pill = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.05), Inches(2.4), Inches(0.32)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = PILL_BG
    pill.line.fill.background()
    pt = pill.text_frame.paragraphs[0]
    pt.alignment = PP_ALIGN.CENTER
    r = pt.add_run()
    r.text = f'For: {audience}'
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for i, b in enumerate(bullets):
        txbox(slide, f'• {b}', 0.5, 1.55 + i * 0.8, 12.3, 0.7, size=17, color=MID)

# ── Slide 1: Title ─────────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
# Red accent bar
bar = s.shapes.add_shape(1, Inches(0.8), Inches(2.65), Inches(11.7), Inches(0.07))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()
txbox(s, 'Transcript Intelligence', 0.8, 1.3, 11.7, 1.1,
      size=46, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txbox(s, 'AegisCloud Meeting Analytics', 0.8, 2.8, 11.7, 0.7,
      size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txbox(s, 'Theme classification and actionable insights from 100 customer meetings',
      0.8, 3.6, 11.7, 0.55, size=16,
      color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

# ── Slide 2: Approach 1 — Rule-Based ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
section_slide(s, 'Approach 1 — Rule-Based Keyword Matching', 'All', [
    'Every theme must be hand-coded — anything unexpected is invisible to the rules',
    '"Pipeline failure" and "ingestion outage" mean the same thing; the rules treat them as different',
    'Each product change or new customer term requires a manual update — breaks silently with no warning',
])

# ── Slide 3: Approach 2 — TF-IDF + KMeans ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
section_slide(s, 'Approach 2 — TF-IDF + KMeans', 'All', [
    'Number of themes must be guessed upfront — wrong K forces unrelated topics together',
    'Bag-of-words misses meaning: "outage remediation" and "post-mortem" treated as unrelated',
    'Theme names are mechanical — top centroid terms like "renewal / competitive / pricing / outage"',
])

# ── Slide 4: Approach 3 — Final ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
section_slide(s, 'Approach 3 — Embeddings + HDBSCAN + LLM', 'All', [
    'Each topic phrase embedded independently — a meeting about HIPAA and API failures maps to two distinct clusters',
    'HDBSCAN found 26 themes from the data — no K required, density-adaptive',
    'LLM names the clusters, it does not form them — the hard work is done by the algorithm',
])

# ── Slide 5: Database Schema ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
txbox(s, 'Database Schema', 0.5, 0.25, 12.3, 0.7, size=28, bold=True, color=DARK)
if SCHEMA_IMG.exists():
    s.shapes.add_picture(clean_img(SCHEMA_IMG), Inches(0.4), Inches(1.0), Inches(7.5), Inches(5.8))
txbox(s, '• 6 base tables loaded from raw JSON — meetings, summaries, key moments, action items, participants, transcripts',
      8.1, 1.5, 4.9, 1.2, size=13, color=MID)
txbox(s, '• 3 semantic tables loaded from clustering — all insight queries join both layers in one schema with no cross-database joins',
      8.1, 3.0, 4.9, 1.4, size=13, color=MID)

# ── Slide 6: 3.1 Meeting Breakdown ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'Meeting Breakdown by Call Type and Product', 'All', [
    '100 meetings: roughly half support, a quarter renewals, a quarter internal',
    'Detect and Comply appear in the most meetings across the dataset',
    'Sets the scale before any analysis begins',
], CHARTS / '00_dataset_overview.png')

# ── Slide 7: 4.1 UMAP Scatter ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'UMAP Scatter — 343 Topic Phrases, 26 Clusters', 'All', [
    'Each dot is a topic phrase from meeting summaries, coloured by cluster',
    '26 themes emerged naturally — no K was specified',
    'Detect phrases cluster tightly together; Comply sits as a separate island',
], CHARTS / '01_umap_scatter.png')

# ── Slide 8: 5.1 Theme Sentiment Heatmap ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'Theme Sentiment Heatmap', 'All Leadership', [
    'Sorted most negative (top) to most positive (bottom)',
    'Detect-related themes dominate the red zone; Comply dominates the green zone',
    'One product is in crisis — the other is the strongest asset in the portfolio',
], CHARTS / '03_theme_sentiment_heatmap.png')

# ── Slide 9: 6.1 Churn Signals by Product ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'Churn Signals by Product', 'Sales, CS', [
    'Churn signals = moments where customers explicitly expressed risk of leaving',
    'Detect generates far more churn signals than any other product',
    'Direct input for account prioritisation and recovery planning',
], CHARTS / '04_churn_by_product.png')

# ── Slide 10: 6.2 Accounts Needing Immediate Follow-Up ────────────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'Accounts Needing Immediate Follow-Up', 'Sales, CS', [
    '38 meetings with negative sentiment and explicit churn signals',
    'Ranked by severity — top rows are the most urgent accounts',
    'Action list, not a trend — these accounts need a call this week',
], CHARTS / '04b_watchlist.png')

# ── Slide 11: 6.3 Which Products Are Generating the Most Risk? ────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'Which Products Are Generating the Most Risk?', 'Engineering, Product', [
    'Three signals per product: total meetings, technical issues, and churn signals',
    'Detect leads on both technical issues and churn — the outage created a long tail',
    'Shows which product needs reliability investment first',
], CHARTS / '05_product_signals.png')

# ── Slide 12: 6.4 Where Customers Are Happy ───────────────────────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'Where Customers Are Happy', 'Product, Marketing, Sales, CS', [
    'Comply has the highest praise density of any product',
    'Comply renewal conversations are the most positive in the dataset',
    'Comply v2 is the good news story to pair with the Detect recovery message',
], CHARTS / '06_positive_signals.png')

# ── Slide 13: 7.1 How the Detect Outage Affected Renewal Calls ────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'How the Detect Outage Affected Renewal Calls', 'Engineering, Sales', [
    'Renewal meetings mentioning Detect are significantly more negative than those that do not',
    'Over half of Detect renewal meetings carry an explicit churn signal',
    'Outage conversations are bleeding into commercial calls that should be about growth',
], CHARTS / '07_detect_external_impact.png')

# ── Slide 14: 7.2 What Customers Are Asking For and How Urgently ──────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'What Customers Are Asking For and How Urgently', 'Product (CPO)', [
    'Same feature request carries different urgency depending on the account situation',
    'Blocked (red) = at-risk customer — treat as P0 and act immediately',
    'Growing (green) = healthy account asking for more — add to the roadmap',
], CHARTS / '08_feature_gaps_by_product.png')

# ── Slide 15: 7.3 Who Owns the Follow-Up Work? ────────────────────────────
s = prs.slides.add_slide(BLANK)
chart_slide(s, 'Who Owns the Follow-Up Work?', 'Operations, Engineering, CS, Sales', [
    'Left: which themes generate the most follow-up actions across all meetings',
    'Engineering carries the heaviest Detect load — reliability slips delay Comply v2',
    'Identifies where capacity bottlenecks sit before they become a problem',
], CHARTS / '09_action_item_owners.png')

# ── Slide 16: Three Things to Act On ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
txbox(s, 'Three Things to Act On', 0.5, 0.3, 12.3, 0.8,
      size=30, bold=True, color=WHITE)
items = [
    ('1', 'Fix Detect reliability before the next QBR',
          'It is the #1 churn driver across the dataset'),
    ('2', 'Call the 38 high-risk accounts this week',
          'They have negative sentiment and explicit churn signals'),
    ('3', 'Lead with Comply v2 for at-risk Detect accounts',
          'It is the strongest counter-narrative in the data'),
]
for i, (num, heading, sub) in enumerate(items):
    top = 1.4 + i * 1.7
    # Number circle
    circ = s.shapes.add_shape(9, Inches(0.5), Inches(top), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    cp = circ.text_frame.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = num
    cr.font.size = Pt(18)
    cr.font.bold = True
    cr.font.color.rgb = WHITE
    txbox(s, heading, 1.25, top - 0.04, 11.5, 0.45, size=18, bold=True, color=WHITE)
    txbox(s, sub,     1.25, top + 0.42, 11.5, 0.38, size=14, color=RGBColor(0xCC, 0xCC, 0xCC))

prs.save(OUT)
print(f'Saved: {OUT}')
